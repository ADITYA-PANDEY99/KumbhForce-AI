import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

prs = Presentation()

# Apply 16:9 widescreen layout
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Dark theme colors (Command Aurora style)
BG_COLOR = RGBColor(9, 10, 15)        # Midnight Blue/Black
PRIMARY_COLOR = RGBColor(6, 182, 212) # Neon Teal
TEXT_LIGHT = RGBColor(243, 244, 246)  # Warm White
TEXT_MUTED = RGBColor(156, 163, 175)  # Muted Gray
ACCENT_GREEN = RGBColor(16, 185, 129) # Vivid Emerald

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    return title_box

# Slide 1: Title
blank_slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide1)

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "KumbhForce AI"
p1.font.name = "Arial"
p1.font.size = Pt(64)
p1.font.bold = True
p1.font.color.rgb = PRIMARY_COLOR

p2 = tf.add_paragraph()
p2.text = "Synchronizing Humanity at Scale"
p2.font.name = "Arial"
p2.font.size = Pt(28)
p2.font.bold = True
p2.font.color.rgb = TEXT_LIGHT
p2.space_before = Pt(15)

p3 = tf.add_paragraph()
p3.text = "Autonomous Volunteer Command & Operations Intelligence Platform for Mahakumbh\nPresented by: Aditya Pandey"
p3.font.name = "Arial"
p3.font.size = Pt(16)
p3.font.color.rgb = TEXT_MUTED
p3.space_before = Pt(30)

# Slide templates for other slides
slides_data = [
    {
        "title": "Operating at the Limits of Human Scale",
        "bullets": [
            "120 Million+ Pilgrims: Population size larger than most European nations, shifting daily.",
            "Dynamic Floodplain Grid: Temporary pontoon bridges and riverbanks shifting dynamically.",
            "Density Crises: Pilgrim density exceeding 4 people/sqm at transition gates and rail hubs."
        ]
    },
    {
        "title": "Why Traditional Management Fails",
        "bullets": [
            "Roster Blindness: Static volunteer shifts with no live location, skill, or fatigue tracking.",
            "Response Latency: Cascade communication or radio delays introduce a 15+ minute response gap.",
            "Reactive Posture: Command centers react to crises only after bottlenecks materialize."
        ]
    },
    {
        "title": "KumbhForce AI: Command Redefined",
        "bullets": [
            "Digital Twin Telemetry: Live spatial tracking of all 12 active sectors of the Mahakumbh floodplain.",
            "Constraint-Aware Solver: Linear program matching volunteers by proximity, skillsets, and fatigue.",
            "Simulation Sandbox: Run Monte Carlo scenarios to stress-test capacity and generate briefings."
        ]
    },
    {
        "title": "Intelligent Operations Pipeline",
        "bullets": [
            "1. Telemetry Ingress: Live pilgrim density metrics, GPS feeds, and incident triggers.",
            "2. AI Decision Engine: Evaluates threat profiles and computes sector priorities.",
            "3. Roster Optimization Solver: Distributes personnel to clear critical gaps and balance workloads.",
            "4. Interactive HUD: Streams live dispatcher alerts and explainable recommendation trails."
        ]
    },
    {
        "title": "Digital Command Center HUD",
        "bullets": [
            "12-Sector Twin Grid: Interactive matrix displaying sector health, crowd volumes, and staff counts.",
            "Live Operations Feed: Real-time dispatcher activity stream tracking dispatches and rerouting.",
            "Explainability Panel: Instant breakdowns of volunteer skill matches, fatigue indexes, and distances."
        ]
    },
    {
        "title": "Workforce Optimizer Engine",
        "bullets": [
            "Skill-Constraint Solver: Automatically allocates specialist volunteers (First Aid, Languages).",
            "Fatigue-Aware Routing: Reassigns standby personnel without overloading active teams.",
            "Auto-Resolve Gaps: Execute system-wide personnel balancing with a single button click."
        ]
    },
    {
        "title": "Scenario Stress Sandbox",
        "bullets": [
            "Interactive Sliders: Stress-test operations by scaling crowd sizes and reducing volunteers.",
            "Capacity Forecast Chart: Real-time graph of risk index and crowd trends over an 8-hour horizon.",
            "Scenario Presets: Instant profiles for 'Crowd Surge', 'Heatwave Emergency', and 'Volunteer Shortage'."
        ]
    },
    {
        "title": "60-Second Incident Response Flow",
        "bullets": [
            "1. Detection: Medical emergency ticket logged at S-02 (Triveni Point) with Critical severity.",
            "2. Profile Search: Solver scans nearby low-fatigue, first-aid certified volunteers.",
            "3. Candidate Proposal: Rohit Kumar (0.4 km away) recommended with 95% confidence score.",
            "4. Dispatch Authorization: Dispatcher executes rebalancing command via one-click CTA.",
            "5. Resolution: Volunteer arrives within 3.8 minutes, beating standard SLA targets."
        ]
    },
    {
        "title": "Operations Decision Assistant (AI Copilot)",
        "bullets": [
            "Natural Language Queries: Command queries resolved instantly with database citations.",
            "SLA Analysis Headers: Confidence score, estimated response time, and impact metrics.",
            "Executable CTAs: Actionable buttons linked directly to dispatches and map views."
        ]
    },
    {
        "title": "Projected Operational Outcomes",
        "bullets": [
            "Placement Rates: Volunteer deployment success increased from 62% to 94%.",
            "Response Latency: Average dispatch latency reduced from 14.5 minutes to 6.2 minutes (-57%).",
            "Safety Gaps: Floodplain safety coverage raised from 71% to 95%.",
            "Crowd Bottlenecks: Projected crowd surge risks decreased by 46% across target checkpoints."
        ]
    },
    {
        "title": "Comparative Advantage Matrix",
        "bullets": [
            "Live 12-Sector Twin vs. spreadsheets: Static rows vs. real-time dynamic visual telemetry.",
            "Constraint Solver vs. basic scheduling: Proximity, skill, and fatigue computed dynamically.",
            "Predictive Sandbox vs. reactive dispatch: Proactive trend graphs vs. late crisis management."
        ]
    },
    {
        "title": "Thank You",
        "bullets": [
            "KumbhForce AI: Autonomous Volunteer Operations Command Center",
            "Lead Creator: Aditya Pandey",
            "Web Demo: kumbhforce-ai.vercel.app",
            "GitHub Repository: github.com/ADITYA-PANDEY99/KumbhForce-AI"
        ]
    }
]

for s_data in slides_data:
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide)
    add_title(slide, s_data["title"])
    
    # Add bullet points box
    bullets_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(s_data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        p.font.name = "Arial"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(20)

output_path = os.path.join(os.getcwd(), "pitch_deck.pptx")
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
