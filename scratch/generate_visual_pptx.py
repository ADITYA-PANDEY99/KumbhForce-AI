import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Style colors (Dark Aurora Theme)
BG_COLOR = RGBColor(9, 10, 15)          # Midnight Blue/Black
CYAN = RGBColor(6, 182, 212)            # Accent Cyan
GREEN = RGBColor(16, 185, 129)          # Accent Green
RED = RGBColor(239, 68, 68)             # Alert Red
TEXT_LIGHT = RGBColor(243, 244, 246)    # Gray 100
TEXT_MUTED = RGBColor(156, 163, 175)    # Gray 400
PANEL_BG = RGBColor(17, 24, 39)         # Dark Gray Panel

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title, category="SYSTEM TELEMETRY"):
    # Category badge
    cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = CYAN
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(11.83), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT

def create_card(slide, left, top, width, height, title, value, desc, border_color=CYAN):
    # Base Panel
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    
    # Content Frame
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = title.upper()
    p_title.font.name = "Arial"
    p_title.font.size = Pt(11)
    p_title.font.bold = True
    p_title.font.color.rgb = border_color
    
    p_val = tf.add_paragraph()
    p_val.text = value
    p_val.font.name = "Arial"
    p_val.font.size = Pt(24)
    p_val.font.bold = True
    p_val.font.color.rgb = TEXT_LIGHT
    p_val.space_before = Pt(8)
    
    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.space_before = Pt(8)

def add_arrow(slide, left, top, width, height, direction="down", color=CYAN):
    # Simulates flows
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW if direction == "down" else MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ==========================================
# SLIDE 1: TITLE (Minimalist, Hero Title)
# ==========================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide1)

# Glowing orbit simulation (simple circles behind text)
orbit = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.66), Inches(0.75), Inches(6.0), Inches(6.0))
orbit.fill.solid()
orbit.fill.fore_color.rgb = RGBColor(12, 20, 35)
orbit.line.color.rgb = CYAN
orbit.line.width = Pt(1)

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
p1.text = "KUMBHFORCE AI"
p1.font.name = "Arial"
p1.font.size = Pt(60)
p1.font.bold = True
p1.font.color.rgb = CYAN

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.text = "Autonomous Volunteer Command & Spatial Intelligence"
p2.font.name = "Arial"
p2.font.size = Pt(22)
p2.font.color.rgb = TEXT_LIGHT
p2.space_before = Pt(15)

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.text = "Mahakumbh Staffing War Room\nDesigned & Built by Aditya Pandey"
p3.font.name = "Arial"
p3.font.size = Pt(13)
p3.font.color.rgb = TEXT_MUTED
p3.space_before = Pt(40)


# ==========================================
# SLIDE 2: SCALE (3 Statistic Cards)
# ==========================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide2)
add_header(slide2, "Scale of the Problem", "Floodplain constraints")

create_card(slide2, Inches(1.0), Inches(2.2), Inches(3.5), Inches(4.2), "Demographics", "120M+ Pilgrims", "Largest human gathering on Earth. Larger population footprint than most European nations transit daily.", CYAN)
create_card(slide2, Inches(4.9), Inches(2.2), Inches(3.5), Inches(4.2), "Topography", "Dynamic Grid", "Floodplain structures, pathways, and pontoon bridges change shape dynamically based on seasonal water levels.", CYAN)
create_card(slide2, Inches(8.8), Inches(2.2), Inches(3.5), Inches(4.2), "SLA Bottleneck", "Density Crisis", "Bathing gates and transit terminals exceed 4 people/sqm, triggering extreme crowd surge risks.", RED)


# ==========================================
# SLIDE 3: TRADITIONAL VS KUMBHFORCE AI (Comparison Lanes)
# ==========================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide3)
add_header(slide3, "Traditional Management vs. KumbhForce AI", "Operations Paradigm Shift")

# Lane 1 Header: Traditional
l1_h = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.4))
l1_h.fill.solid()
l1_h.fill.fore_color.rgb = RED
l1_h.line.fill.background()
tf_l1 = l1_h.text_frame
tf_l1.paragraphs[0].text = "TRADITIONAL REACTIVE LOGISTICS"
tf_l1.paragraphs[0].font.size = Pt(11)
tf_l1.paragraphs[0].font.bold = True
tf_l1.paragraphs[0].font.color.rgb = TEXT_LIGHT

# Lane 1 Flow Blocks
create_card(slide3, Inches(1.0), Inches(2.6), Inches(3.5), Inches(1.6), "Roster Blindness", "Static Shifts", "No tracking of fatigue, location, or active responder skill levels.", RED)
add_arrow(slide3, Inches(4.6), Inches(3.2), Inches(0.2), Inches(0.4), "right", RED)
create_card(slide3, Inches(4.9), Inches(2.6), Inches(3.5), Inches(1.6), "Communication Gap", "Radio Delays", "15+ minute dispatch gaps via voice calls and cascading radio trees.", RED)
add_arrow(slide3, Inches(8.5), Inches(3.2), Inches(0.2), Inches(0.4), "right", RED)
create_card(slide3, Inches(8.8), Inches(2.6), Inches(3.5), Inches(1.6), "Crisis Response", "Reactive Mode", "Personnel deployed only after safety bottlenecks or medical emergencies materialize.", RED)

# Lane 2 Header: KumbhForce AI
l2_h = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.4))
l2_h.fill.solid()
l2_h.fill.fore_color.rgb = GREEN
l2_h.line.fill.background()
tf_l2 = l2_h.text_frame
tf_l2.paragraphs[0].text = "KUMBHFORCE AI PROACTIVE INTELLIGENCE"
tf_l2.paragraphs[0].font.size = Pt(11)
tf_l2.paragraphs[0].font.bold = True
tf_l2.paragraphs[0].font.color.rgb = TEXT_LIGHT

# Lane 2 Flow Blocks
create_card(slide3, Inches(1.0), Inches(5.2), Inches(3.5), Inches(1.6), "Unified HUD", "Digital Twin", "Live spatial visualization of 12 sectors with active telemetry.", GREEN)
add_arrow(slide3, Inches(4.6), Inches(5.8), Inches(0.2), Inches(0.4), "right", GREEN)
create_card(slide3, Inches(4.9), Inches(5.2), Inches(3.5), Inches(1.6), "Roster Optimization", "Linear Program", "Constraint matching solver balances proximity, skills, and fatigue.", GREEN)
add_arrow(slide3, Inches(8.5), Inches(5.8), Inches(0.2), Inches(0.4), "right", GREEN)
create_card(slide3, Inches(8.8), Inches(5.2), Inches(3.5), Inches(1.6), "Risk Mitigation", "Predictive Prevention", "Simulation sandbox tracks risk indices to deploy help before issues rise.", GREEN)


# ==========================================
# SLIDE 4: ARCHITECTURE (Vertical Flowchart)
# ==========================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide4)
add_header(slide4, "Operations System Architecture Pipeline", "Data ingestion & optimization flow")

steps = [
    ("1. Telemetry Ingress", "GPS coordinates, pilgrim density sensors, and incident triggers"),
    ("2. AI Decision Layer", "Threat score calculator and sector priority ranker"),
    ("3. Optimization Engine", "Constraint matching solver mapping skills, distance, and fatigue"),
    ("4. Digital Twin Center", "Unified HUD visualizer mapping 12-sector health matrices")
]

for idx, (title, desc) in enumerate(steps):
    top_offset = Inches(2.0 + (idx * 1.25))
    
    # Step Card
    shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), top_offset, Inches(10.33), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_BG
    shape.line.color.rgb = CYAN
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{title.upper()}  ·  {desc}"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    
    # Connecting Arrow
    if idx < 3:
        add_arrow(slide4, Inches(6.4), top_offset + Inches(0.8), Inches(0.5), Inches(0.45), "down", CYAN)


# ==========================================
# SLIDE 5: COMMAND CENTER SCREEN (Annotated HUD)
# ==========================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide5)
add_header(slide5, "Digital Command Center HUD", "Live Spatial Twin Matrix")

# Screenshot Frame
frame = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.2), Inches(7.5), Inches(4.2))
frame.fill.solid()
frame.fill.fore_color.rgb = RGBColor(12, 20, 35)
frame.line.color.rgb = CYAN
frame.line.width = Pt(2)
tf_fr = frame.text_frame
tf_fr.paragraphs[0].text = "[ UI SCREENSHOT: DIGITAL TWIN 12-SECTOR GRID ]\n\nShows active telemetry, volunteer coverage indexes, explainability breakdowns, and category-filtered dispatcher logs."
tf_fr.paragraphs[0].alignment = PP_ALIGN.CENTER
tf_fr.paragraphs[0].font.size = Pt(14)
tf_fr.paragraphs[0].font.color.rgb = TEXT_MUTED

# Callout Cards
create_card(slide5, Inches(9.0), Inches(2.2), Inches(3.3), Inches(1.2), "Component 1", "Live Telemetry", "Streams GPS updates & density index.", CYAN)
create_card(slide5, Inches(9.0), Inches(3.7), Inches(3.3), Inches(1.2), "Component 2", "Sector Intel", "Explainable skill/fatigue matrix.", CYAN)
create_card(slide5, Inches(9.0), Inches(5.2), Inches(3.3), Inches(1.2), "Component 3", "Incident Monitor", "Triage ticketing and routing logs.", GREEN)


# ==========================================
# SLIDE 6: WORKFORCE OPTIMIZER SCREEN (Visual Pipeline)
# ==========================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide6)
add_header(slide6, "Workforce Optimizer Engine", "Real-Time Constraint Dispatcher")

# Screenshot Frame
frame = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.2), Inches(7.5), Inches(4.2))
frame.fill.solid()
frame.fill.fore_color.rgb = RGBColor(12, 20, 35)
frame.line.color.rgb = CYAN
frame.line.width = Pt(2)
tf_fr = frame.text_frame
tf_fr.paragraphs[0].text = "[ UI SCREENSHOT: WORKFORCE OPTIMIZER ]\n\nContrasts volunteer profiles, distance thresholds, and fatigue indexes. Includes single-click Auto-Resolve buttons."
tf_fr.paragraphs[0].alignment = PP_ALIGN.CENTER
tf_fr.paragraphs[0].font.size = Pt(14)
tf_fr.paragraphs[0].font.color.rgb = TEXT_MUTED

# Vertical Pipeline Flow
pipes = [
    ("Skills Match", "First aid / Languages"),
    ("Proximity", "Physical distance math"),
    ("Fatigue Cap", "Duty cycle indexes"),
    ("Optimal Match", "Executed dispatch")
]

for idx, (title, desc) in enumerate(pipes):
    top_offset = Inches(2.2 + (idx * 1.05))
    card = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), top_offset, Inches(3.3), Inches(0.8))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL_BG
    card.line.color.rgb = CYAN if idx < 3 else GREEN
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    p_title = tf.paragraphs[0]
    p_title.text = f"{idx+1}. {title.upper()}"
    p_title.font.size = Pt(11)
    p_title.font.bold = True
    p_title.font.color.rgb = CYAN if idx < 3 else GREEN
    
    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(10)
    p_desc.font.color.rgb = TEXT_LIGHT


# ==========================================
# SLIDE 7: SCENARIO SIMULATOR SCREEN (Stress-Test Sandbox)
# ==========================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide7)
add_header(slide7, "Scenario Stress Simulator Sandbox", "What-If Forecasting Dashboard")

# Screenshot Frame
frame = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.2), Inches(7.5), Inches(4.2))
frame.fill.solid()
frame.fill.fore_color.rgb = RGBColor(12, 20, 35)
frame.line.color.rgb = CYAN
frame.line.width = Pt(2)
tf_fr = frame.text_frame
tf_fr.paragraphs[0].text = "[ UI SCREENSHOT: SCENARIO STRESS SIMULATOR ]\n\nShows capacity trend curves, risk dials, preset configurations (Crowd Surge, Dehydration Wave), and slider controls."
tf_fr.paragraphs[0].alignment = PP_ALIGN.CENTER
tf_fr.paragraphs[0].font.size = Pt(14)
tf_fr.paragraphs[0].font.color.rgb = TEXT_MUTED

# Steps Pipeline
create_card(slide7, Inches(9.0), Inches(2.2), Inches(3.3), Inches(1.2), "Step 1", "Before State", "Stable volunteer metrics within normal SLA bounds.", CYAN)
create_card(slide7, Inches(9.0), Inches(3.7), Inches(3.3), Inches(1.2), "Step 2", "Surge Event", "Adjust sliders to simulate crowd surges.", RED)
create_card(slide7, Inches(9.0), Inches(5.2), Inches(3.3), Inches(1.2), "Step 3", "Optimized State", "System recalculates routes to restore safety.", GREEN)


# ==========================================
# SLIDE 8: PREDICTIVE STAFFING SCREEN (Forecast Timeline)
# ==========================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide8)
add_header(slide8, "Predictive Staffing Engine", "Crowd Density Forecast Grid")

# Timeline cards
times = [
    ("NOW", "Current grid capacity: 487 active, 4 incidents reported"),
    ("1H FORECAST", "Sangam Ghat expected crowd +13%. Gap: -22 vols."),
    ("3H FORECAST", "Triveni Point reaches Critical density. Gap: -18 vols."),
    ("6H FORECAST", "Peak gap projected: -58 volunteers in S-01 sector."),
    ("12H FORECAST", "Inflow rates decline; standard monitoring resumes.")
]

for idx, (time_lbl, desc) in enumerate(times):
    left_offset = Inches(0.75 + (idx * 2.4))
    
    # Timeline Node
    card = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_offset, Inches(2.6), Inches(2.2), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL_BG
    card.line.color.rgb = RED if idx == 3 else CYAN
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    p_title = tf.paragraphs[0]
    p_title.text = time_lbl
    p_title.font.name = "Arial"
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = RED if idx == 3 else CYAN
    
    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.space_before = Pt(15)
    
    # Arrow
    if idx < 4:
        add_arrow(slide8, left_offset + Inches(2.2), Inches(4.1), Inches(0.2), Inches(0.4), "right", CYAN)


# ==========================================
# SLIDE 9: AI COPILOT (Interactive Conversation Flow)
# ==========================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide9)
add_header(slide9, "AI Operations Copilot Interface", "Conversational CLI & Clickable Executions")

# Chat Container
chat_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.2), Inches(7.5), Inches(4.2))
chat_box.fill.solid()
chat_box.fill.fore_color.rgb = PANEL_BG
chat_box.line.color.rgb = CYAN
chat_box.line.width = Pt(1.5)

tf_chat = chat_box.text_frame
tf_chat.word_wrap = True

p_q = tf_chat.paragraphs[0]
p_q.text = "USER: \"How is the crowd density at Sangam Ghat S-01, and how do we resolve the current staff gap?\""
p_q.font.name = "Arial"
p_q.font.size = Pt(13)
p_q.font.bold = True
p_q.font.color.rgb = TEXT_LIGHT

p_a = tf_chat.add_paragraph()
p_a.text = "COPILOT: \"S-01 Sangam Ghat is projected to experience a density surge of +13% in 60 minutes. Recommended action: Redeploy 12 volunteers from S-03 Ram Ghat. Confidence: 95%.\""
p_a.font.name = "Arial"
p_a.font.size = Pt(13)
p_a.font.color.rgb = GREEN
p_a.space_before = Pt(20)

p_cta = tf_chat.add_paragraph()
p_cta.text = "[ BUTTON: EXECUTE VOLUNTEER DISPATCH REALLOCATION ]"
p_cta.font.name = "Arial"
p_cta.font.size = Pt(11)
p_cta.font.bold = True
p_cta.font.color.rgb = CYAN
p_cta.space_before = Pt(30)

# Annotation Cards
create_card(slide9, Inches(9.0), Inches(2.2), Inches(3.3), Inches(1.8), "Core Feature", "SLA Analytics", "Every query returns Confidence Scores, SLA impact ratings, and data citations.", CYAN)
create_card(slide9, Inches(9.0), Inches(4.6), Inches(3.3), Inches(1.8), "Control Logic", "Clickable CTAs", "Operations chatbot triggers dispatches, redirects pages, or runs stress simulations.", GREEN)


# ==========================================
# SLIDE 10: INCIDENT RESPONSE JOURNEY (Process Flow)
# ==========================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide10)
add_header(slide10, "Incident Response Journey Flow", "Telemetry trigger to resolution track")

journey = [
    ("1. Alert", "Critical triage ticket generated at S-02."),
    ("2. Analysis", "Threat index assessed, SLA target: 4.5m."),
    ("3. Match", "Solver selects रोहित कुमार (First aid certified)."),
    ("4. Dispatch", "One-click authorization routing command."),
    ("5. Done", "Incident resolved within 3.8 minutes.")
]

for idx, (step_lbl, desc) in enumerate(journey):
    left_offset = Inches(0.75 + (idx * 2.4))
    
    # Process Node
    card = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_offset, Inches(2.6), Inches(2.2), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL_BG
    card.line.color.rgb = GREEN if idx == 4 else CYAN
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    p_title = tf.paragraphs[0]
    p_title.text = step_lbl
    p_title.font.name = "Arial"
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = GREEN if idx == 4 else CYAN
    
    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.space_before = Pt(15)
    
    # Arrow
    if idx < 4:
        add_arrow(slide10, left_offset + Inches(2.2), Inches(4.1), Inches(0.2), Inches(0.4), "right", CYAN)


# ==========================================
# SLIDE 11: DEPLOYMENT & FULL STACK (Infrastructure Architecture)
# ==========================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide11)
add_header(slide11, "Infrastructure & Live Production Stack", "Multi-Cloud deployment architecture")

# 4 Columns for stack layers
layers = [
    ("FRONTEND", "Vercel Cloud", "Next.js 15 app router\nTailwind CSS styling\nHeader status badges\nOffline fallbacks", CYAN),
    ("BACKEND", "Render Cloud", "FastAPI Python ASGI\nUvicorn startup engine\nAuto-deploy hooks\nOfficial /health monitor", CYAN),
    ("DATABASE", "SQLite Layer", "Auto schema generator\nSQLAlchemy ORM\nIncident metrics tables\nUser registry store", CYAN),
    ("INTEGRATIONS", "Operations Platform", "Swagger API docs\nGitHub VCS repository\nFull stack connected\nDemonstrable live HUD", GREEN)
]

for idx, (title, service, details, color) in enumerate(layers):
    left_offset = Inches(0.75 + (idx * 3.0))
    
    card = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_offset, Inches(2.2), Inches(2.7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL_BG
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    tf = card.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(11)
    p_title.font.bold = True
    p_title.font.color.rgb = color
    
    p_srv = tf.add_paragraph()
    p_srv.text = service
    p_srv.font.name = "Arial"
    p_srv.font.size = Pt(20)
    p_srv.font.bold = True
    p_srv.font.color.rgb = TEXT_LIGHT
    p_srv.space_before = Pt(10)
    
    p_det = tf.add_paragraph()
    p_det.text = details
    p_det.font.name = "Arial"
    p_det.font.size = Pt(11)
    p_det.font.color.rgb = TEXT_MUTED
    p_det.space_before = Pt(15)


# ==========================================
# SLIDE 12: CAPABILITIES (6-Grid Checklist)
# ==========================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide12)
add_header(slide12, "Platform Capabilities", "Core functional dimensions")

caps = [
    ("Predictive Staffing", "Anticipate crowd volume anomalies & volunteer shortages."),
    ("Digital Command Twin", "12-sector matrix visual telemetry HUD."),
    ("AI Copilot Assistant", "Conversational CLI decision triggers & citations."),
    ("Workforce Optimization", "Constraint solver routing proximity & fatigue."),
    ("Scenario Simulator", "Stress-test sandbox with presets & SVG trend charts."),
    ("Full Stack Integration", "Live REST API server with client offline fallbacks.")
]

for idx, (title, desc) in enumerate(caps):
    row = idx // 3
    col = idx % 3
    
    left_offset = Inches(0.75 + (col * 4.0))
    top_offset = Inches(2.2 + (row * 2.3))
    
    create_card(slide12, left_offset, top_offset, Inches(3.7), Inches(2.0), "✓ " + title, "Capability", desc, GREEN)


# ==========================================
# SLIDE 13: THANK YOU & OUTRO (Widescreen QR Card)
# ==========================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(slide13)
add_header(slide13, "Synchronizing Humanity at Scale", "Contact & Repository Credentials")

# QR Code Box Placeholder
qr_box = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.2), Inches(4.5), Inches(4.2))
qr_box.fill.solid()
qr_box.fill.fore_color.rgb = RGBColor(12, 20, 35)
qr_box.line.color.rgb = CYAN
qr_box.line.width = Pt(1.5)
tf_qr = qr_box.text_frame
tf_qr.paragraphs[0].text = "[ QR CODE MOCKUP ]\n\nScan to access live platform deployment and GitHub code files."
tf_qr.paragraphs[0].alignment = PP_ALIGN.CENTER
tf_qr.paragraphs[0].font.size = Pt(13)
tf_qr.paragraphs[0].font.color.rgb = TEXT_MUTED

# Links Card
link_box = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(2.2), Inches(6.33), Inches(4.2))
link_box.fill.solid()
link_box.fill.fore_color.rgb = PANEL_BG
link_box.line.color.rgb = CYAN
link_box.line.width = Pt(1.5)

tf_lnk = link_box.text_frame
tf_lnk.word_wrap = True

p_demo = tf_lnk.paragraphs[0]
p_demo.text = "DEMO PLATFORM URL:"
p_demo.font.size = Pt(11)
p_demo.font.bold = True
p_demo.font.color.rgb = CYAN

p_demolink = tf_lnk.add_paragraph()
p_demolink.text = "kumbhforce-ai.vercel.app"
p_demolink.font.size = Pt(18)
p_demolink.font.bold = True
p_demolink.font.color.rgb = TEXT_LIGHT
p_demolink.space_before = Pt(4)

p_git = tf_lnk.add_paragraph()
p_git.text = "GITHUB REPOSITORY:"
p_git.font.size = Pt(11)
p_git.font.bold = True
p_git.font.color.rgb = CYAN
p_git.space_before = Pt(15)

p_gitlink = tf_lnk.add_paragraph()
p_gitlink.text = "github.com/ADITYA-PANDEY99/KumbhForce-AI"
p_gitlink.font.size = Pt(18)
p_gitlink.font.bold = True
p_gitlink.font.color.rgb = TEXT_LIGHT
p_gitlink.space_before = Pt(4)

p_c = tf_lnk.add_paragraph()
p_c.text = "CREATOR DETAILS:"
p_c.font.size = Pt(11)
p_c.font.bold = True
p_c.font.color.rgb = CYAN
p_c.space_before = Pt(15)

p_clink = tf_lnk.add_paragraph()
p_clink.text = "Aditya Pandey  ·  aditya@kumbhforce.ai"
p_clink.font.size = Pt(14)
p_clink.font.bold = True
p_clink.font.color.rgb = GREEN
p_clink.space_before = Pt(4)


# Save
output_path = os.path.join(os.getcwd(), "pitch_deck.pptx")
prs.save(output_path)
print(f"Redesigned presentation saved successfully to {output_path}")
